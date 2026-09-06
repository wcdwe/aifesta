"""
연금 Agent 과제 - 문서 포맷별 추출기 (PDF / DOCX / XLSX / PPTX)

모든 추출기는 동일한 스키마를 반환한다:
    tables: [{"page": int, "table_index": int, "rows": int, "cols": int, "data": [[str, ...], ...]}]
    texts:  [{"page": int, "text": str}]

"page"는 포맷에 따라 의미가 다르다 (PDF=페이지, DOCX=1로 고정된 단일 블록,
XLSX=시트, PPTX=슬라이드 번호). 근거 문서 표시(source_doc + page)에는 이 정의로 충분하다.

extract_any(path)가 확장자를 보고 알맞은 함수로 위임한다.
"""

import os
import sys

import pdfplumber

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import pdf_words  # noqa: E402  (import만으로 Page.chars 전역 패치가 걸린다 - pdf_words.py 참고)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def extract_pdf_tables(pdf_path, page_range=None):
    """PDF에서 페이지별 표를 추출한다."""
    results = []
    with pdfplumber.open(pdf_path) as pdf:
        pages = pdf.pages
        start = 1
        if page_range:
            start = page_range[0]
            pages = pdf.pages[page_range[0] - 1: page_range[1]]

        for i, page in enumerate(pages):
            page_num = start + i
            tables = page.extract_tables()
            for t_idx, table in enumerate(tables):
                cleaned = [
                    [cell.strip() if cell else "" for cell in row]
                    for row in table
                ]
                results.append({
                    "page": page_num,
                    "table_index": t_idx,
                    "rows": len(cleaned),
                    "cols": len(cleaned[0]) if cleaned else 0,
                    "data": cleaned,
                })
    return results


def extract_pdf_text(pdf_path, page_range=None):
    """서술형 텍스트를 페이지 단위로 추출한다."""
    documents = []
    with pdfplumber.open(pdf_path) as pdf:
        pages = pdf.pages
        start = 1
        if page_range:
            start = page_range[0]
            pages = pdf.pages[page_range[0] - 1: page_range[1]]

        for i, page in enumerate(pages):
            page_num = start + i
            text = page.extract_text()
            if text:
                documents.append({"page": page_num, "text": text})
    return documents


def extract_pdf(path):
    return extract_pdf_tables(path), extract_pdf_text(path)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def extract_docx(path):
    import docx

    doc = docx.Document(path)

    tables = []
    for t_idx, table in enumerate(doc.tables):
        data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append({
            "page": 1,
            "table_index": t_idx,
            "rows": len(data),
            "cols": len(data[0]) if data else 0,
            "data": data,
        })

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    texts = [{"page": 1, "text": text}] if text else []

    return tables, texts


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

def extract_xlsx(path):
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)

    tables = []
    texts = []
    for sheet_idx, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if all(c is None for c in row):
                continue
            rows.append(["" if c is None else str(c) for c in row])

        if not rows:
            continue

        tables.append({
            "page": sheet_idx,
            "table_index": 0,
            "rows": len(rows),
            "cols": len(rows[0]) if rows else 0,
            "sheet_name": sheet_name,
            "data": rows,
        })

        # 행 단위 Q&A/레코드 시트(예: FAQ)는 한 행을 한 문단으로 렌더링해서
        # chunk_text.py가 행 단위로 청크를 나눌 수 있게 한다.
        header = rows[0]
        blocks = []
        for row in rows[1:]:
            pairs = [
                f"{h}: {v}" for h, v in zip(header, row) if h and v
            ]
            if pairs:
                blocks.append(" | ".join(pairs))
        sheet_text = "\n\n".join(blocks)
        if sheet_text:
            texts.append({"page": sheet_idx, "text": sheet_text})

    return tables, texts


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)

    tables = []
    texts = []
    for slide_idx, slide in enumerate(list(prs.slides), start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_text_parts.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                tbl = shape.table
                data = [
                    [cell.text.strip() for cell in row.cells]
                    for row in tbl.rows
                ]
                tables.append({
                    "page": slide_idx,
                    "table_index": len(tables),
                    "rows": len(data),
                    "cols": len(data[0]) if data else 0,
                    "data": data,
                })

        slide_text = "\n\n".join(slide_text_parts)
        if slide_text:
            texts.append({"page": slide_idx, "text": slide_text})

    return tables, texts


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


def extract_any(path):
    """확장자를 보고 알맞은 추출기로 위임한다. (tables, texts) 반환."""
    ext = os.path.splitext(path)[1].lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"지원하지 않는 확장자: {ext} ({path})")
    return EXTRACTORS[ext](path)
